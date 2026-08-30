"""
File parser utility.
Extracts plain text from PDF, DOCX, PPTX, TXT, CSV, and XLSX uploads.
"""

import io
import pandas as pd


def parse_uploaded_file(uploaded_file) -> str:
    """
    Accept a Streamlit UploadedFile object and return extracted text.
    Returns an empty string on failure.
    """
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()

    try:
        if name.endswith(".pdf"):
            return _parse_pdf(raw)
        elif name.endswith(".docx"):
            return _parse_docx(raw)
        elif name.endswith(".pptx"):
            return _parse_pptx(raw)
        elif name.endswith(".txt") or name.endswith(".md"):
            return raw.decode("utf-8", errors="ignore")
        elif name.endswith(".csv"):
            return _parse_csv(raw)
        elif name.endswith(".xlsx") or name.endswith(".xls"):
            return _parse_xlsx(raw)
        else:
            return raw.decode("utf-8", errors="ignore")
    except Exception as exc:
        return f"[PARSE ERROR for {uploaded_file.name}: {exc}]"


# ── PDF ──────────────────────────────────────────────────────────────────────

def _parse_pdf(raw: bytes) -> str:
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"[Page {i}]\n{page_text}")
        return "\n\n".join(text_parts)
    except ImportError:
        pass

    # Fallback to PyPDF2
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(raw))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            t = page.extract_text() or ""
            if t.strip():
                parts.append(f"[Page {i}]\n{t}")
        return "\n\n".join(parts)
    except Exception as exc:
        return f"[PDF parse error: {exc}]"


# ── DOCX ─────────────────────────────────────────────────────────────────────

def _parse_docx(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _parse_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


# ── CSV ──────────────────────────────────────────────────────────────────────

def _parse_csv(raw: bytes) -> str:
    df = pd.read_csv(io.BytesIO(raw))
    return df.to_string(index=False)


# ── XLSX ─────────────────────────────────────────────────────────────────────

def _parse_xlsx(raw: bytes) -> str:
    xl = pd.ExcelFile(io.BytesIO(raw))
    parts = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        parts.append(f"[Sheet: {sheet}]\n{df.to_string(index=False)}")
    return "\n\n".join(parts)
